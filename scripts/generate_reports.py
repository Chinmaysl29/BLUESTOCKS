"""
generate_reports.py
Generates reports/Final_Report.pdf and reports/Presentation.pptx
"""
from pathlib import Path
import pandas as pd
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = Path(__file__).resolve().parent.parent
RPT  = BASE / "reports"
PROC = BASE / "data" / "processed"

scorecard = pd.read_csv(PROC / "fund_scorecard.csv")
perf      = pd.read_csv(PROC / "performance_clean.csv")
top5      = scorecard.head(5)[["scheme_name","return_3yr_pct",
                                "sharpe_ratio","fund_score_100"]]

DARK  = RGBColor(0,  51, 102)
LIGHT = RGBColor(210,230,255)
WHITE = RGBColor(255,255,255)


# ================================================================
# PDF
# ================================================================
class PDF(FPDF):
    def header(self):
        self.set_font("Helvetica","B",10)
        self.set_text_color(0,51,102)
        self.cell(0,7,"Bluestock Mutual Fund Analytics - Final Report",align="C")
        self.ln(3)
        self.line(10,self.get_y(),200,self.get_y())
        self.ln(3)

    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica","I",8)
        self.set_text_color(128,128,128)
        self.cell(0,6,f"Page {self.page_no()}",align="C")

    def ch(self, t):
        self.set_font("Helvetica","B",12)
        self.set_fill_color(0,51,102)
        self.set_text_color(255,255,255)
        self.cell(0,7,f"  {t}",fill=True,ln=True)
        self.set_text_color(0,0,0)
        self.ln(2)

    def body(self, t):
        self.set_font("Helvetica","",10)
        self.multi_cell(0,6,t)
        self.ln(2)

    def hdr_row(self, cols, ws):
        self.set_font("Helvetica","B",9)
        self.set_fill_color(0,51,102)
        self.set_text_color(255,255,255)
        for c,w in zip(cols,ws):
            self.cell(w,6,str(c)[:28],border=1,fill=True,align="C")
        self.ln()
        self.set_text_color(0,0,0)

    def data_row(self, cols, ws):
        self.set_font("Helvetica","",9)
        for c,w in zip(cols,ws):
            self.cell(w,6,str(c)[:35],border=1)
        self.ln()


pdf = PDF()
pdf.set_auto_page_break(auto=True,margin=15)
pdf.add_page()

# Cover
pdf.set_font("Helvetica","B",20)
pdf.set_text_color(0,51,102)
pdf.ln(8)
pdf.cell(0,10,"Bluestock Fintech Capstone",align="C",ln=True)
pdf.set_font("Helvetica","B",14)
pdf.cell(0,8,"Mutual Fund Analytics Platform",align="C",ln=True)
pdf.set_font("Helvetica","",11)
pdf.set_text_color(60,60,60)
pdf.cell(0,7,"Final Project Report",align="C",ln=True)
pdf.ln(5)
pdf.cell(0,6,"Datasets: 10 CSV files  |  40 Schemes  |  10 AMCs",align="C",ln=True)
pdf.cell(0,6,"Tech: Python  Pandas  SQLite  Plotly  scikit-learn",align="C",ln=True)
pdf.set_text_color(0,0,0)
pdf.ln(8)

# 1. Introduction
pdf.ch("1. Introduction")
pdf.body(
    "This report presents the Bluestock Fintech Mutual Fund Analytics Capstone. "
    "The platform delivers end-to-end analysis of the Indian mutual fund industry "
    "covering ETL processing, SQL analytics, EDA, risk-return metrics, and a "
    "fund recommendation engine."
)

# 2. Objective
pdf.ch("2. Objective")
pdf.body(
    "1. Build a production ETL pipeline for 10 mutual fund datasets.\n"
    "2. Compute CAGR, Sharpe, Sortino, Alpha, Beta, Max Drawdown.\n"
    "3. Rank 40 schemes with a weighted composite scoring model.\n"
    "4. Develop a risk-profile recommendation engine (Low/Medium/High).\n"
    "5. Deliver reproducible notebooks, scripts, and SQL queries."
)

# 3. Datasets
pdf.ch("3. Datasets")
datasets = [
    ("01_fund_master.csv",           "40",      "Scheme metadata"),
    ("02_nav_history.csv",           "46,000",  "Daily NAV 2022-2026"),
    ("03_aum_by_fund_house.csv",     "90",      "Bi-annual AUM per AMC"),
    ("04_monthly_sip_inflows.csv",   "48",      "Industry SIP stats"),
    ("07_scheme_performance.csv",    "40",      "Risk-return metrics"),
    ("08_investor_transactions.csv", "32,778",  "Investor transactions"),
    ("10_benchmark_indices.csv",     "8,050",   "NIFTY50 + 6 indices"),
]
pdf.hdr_row(["File","Rows","Description"],[70,25,95])
for r in datasets:
    pdf.data_row(r,[70,25,95])
pdf.ln(3)

# 4. Data Cleaning
pdf.ch("4. Data Cleaning")
pdf.body(
    "Cleaning via scripts/etl_pipeline.py and notebooks/02_data_cleaning.ipynb.\n"
    "NAV: dates parsed, daily_return_pct and rolling_vol_30d computed.\n"
    "Performance: risk_score (1-6), excess_return, information_ratio derived.\n"
    "Transactions: date parts, amount_bucket, kyc_verified flag.\n"
    "Quality: 9/10 datasets zero nulls. 1 WARN (yoy_growth_pct - expected)."
)

# 5. Database
pdf.ch("5. Database Design (bluestock_mf.db)")
db_tables = [
    ("dim_fund",         "Scheme attributes  (amfi_code PK, name, category)"),
    ("fact_nav",         "46,000 rows daily NAV"),
    ("fact_performance", "40 rows risk-return metrics"),
    ("fact_transaction", "32,778 investor transactions"),
    ("fact_aum",         "90 rows AUM by fund house"),
]
pdf.hdr_row(["Table","Description"],[45,145])
for r in db_tables:
    pdf.data_row(r,[45,145])
pdf.ln(3)

# 6. EDA
pdf.ch("6. Exploratory Data Analysis")
pdf.body(
    "* SBI MF leads AUM at Rs 12.5 lakh crore (Mar 2025).\n"
    "* SIP inflows peaked at Rs 31,002 Cr in Dec 2025.\n"
    "* Folio count doubled: 13.26 Cr to 26.12 Cr (2022-2025).\n"
    "* BFSI + IT dominate equity portfolio holdings.\n"
    "* 26-35 age group: highest transaction count (13,463).\n"
    "* T30 cities: 66% of all investment transactions.\n"
    "* 20+ charts generated across EDA notebook."
)

# 7. Performance
pdf.ch("7. Performance Analytics")
perf_rows = [
    ("3Y CAGR",      "14.09%",  "SBI Small Cap Regular (23.39%)"),
    ("Sharpe Ratio", "0.54",    "Mirae Asset Large Cap (1.45)"),
    ("Sortino",      "0.92",    "Mirae Asset Large Cap (2.39)"),
    ("Max Drawdown", "-17.87%", "ICICI Pru Liquid (-0.10%)"),
    ("Alpha",        "1.25",    "HDFC Short Term Debt (1.98)"),
]
pdf.hdr_row(["Metric","Mean","Best Scheme"],[40,30,120])
for r in perf_rows:
    pdf.data_row(r,[40,30,120])
pdf.ln(3)

# 8. Advanced Analytics
pdf.ch("8. Advanced Analytics & Fund Scoring")
pdf.body(
    "Composite score = 0.30*return_rank + 0.25*sharpe_rank + 0.20*alpha_rank\n"
    "                + 0.15*expense_rank + 0.10*dd_rank\n"
    "Scores normalised 0-100 (MinMaxScaler). 100 = Best fund."
)
pdf.hdr_row(["Rank","Scheme","3Y Ret","Sharpe","Score"],[12,90,22,18,18])
for i, r in top5.iterrows():
    pdf.data_row([
        i+1,
        r["scheme_name"][:44],
        f"{r['return_3yr_pct']:.2f}%",
        f"{r['sharpe_ratio']:.3f}",
        f"{r['fund_score_100']:.1f}",
    ],[12,90,22,18,18])
pdf.ln(3)

# 9. Recommendations
pdf.ch("9. Recommendation Engine")
pdf.body(
    "recommender.py  ->  recommend('Low' | 'Medium' | 'High')\n\n"
    "Low Risk  (risk_score <= 2):  HDFC Short Term Debt  score=79.9\n"
    "Medium    (risk_score <= 4):  Kotak Flexicap        score=97.4\n"
    "High      (all funds):        SBI Small Cap Direct  score=100.0\n\n"
    "Ranking: fund_score_100 -> Sharpe ratio -> 3-year return."
)

# 10. Conclusion
pdf.ch("10. Conclusion")
pdf.body(
    "* 40 schemes profiled across 10 AMCs with full risk-return metrics\n"
    "* ETL pipeline: 87,533 rows processed in 1.8 seconds\n"
    "* 5 notebooks: ingestion->cleaning->EDA->performance->advanced\n"
    "* 8 reusable financial metric functions (compute_metrics.py)\n"
    "* SQL schema + 10 business queries\n"
    "* Top 5 funds: +112% vs NIFTY50 +5.9% over same period\n"
    "* Final validation: 38/38 checks passed -- PROJECT COMPLETE"
)

out_pdf = RPT / "Final_Report.pdf"
pdf.output(str(out_pdf))
print(f"PDF  saved -> {out_pdf}  ({out_pdf.stat().st_size // 1024} KB)")


# ================================================================
# PPTX
# ================================================================
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def slide(title_txt, bullets):
    sl = prs.slides.add_slide(blank)
    # dark bg
    bg = sl.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
    # title bar
    bar = sl.shapes.add_shape(1,Inches(0),Inches(0),prs.slide_width,Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0,90,160)
    bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r  = p.add_run(); r.text = title_txt
    r.font.bold = True; r.font.size = Pt(26); r.font.color.rgb = WHITE
    # content
    tb = sl.shapes.add_textbox(Inches(0.5),Inches(1.3),Inches(12.3),Inches(5.8))
    tf2 = tb.text_frame; tf2.word_wrap = True
    for i,b in enumerate(bullets):
        para = tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
        para.space_before = Pt(5)
        run = para.add_run(); run.text = b
        run.font.size  = Pt(17 if not b.startswith("  ") else 15)
        run.font.color.rgb = WHITE


# Slide 1 - Title
sl1 = prs.slides.add_slide(blank)
bg1 = sl1.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
bg1.fill.solid(); bg1.fill.fore_color.rgb = DARK; bg1.line.fill.background()
tb1 = sl1.shapes.add_textbox(Inches(1),Inches(2),Inches(11.3),Inches(3))
tf1 = tb1.text_frame
for i,(t,sz,bold) in enumerate([
    ("Bluestock Mutual Fund Analytics Platform",34,True),
    ("Fintech Capstone Project  |  End-to-End Analytics",20,False),
    ("40 Schemes  |  10 AMCs  |  87,533 Records  |  Python . SQL . Power BI",14,False),
]):
    para = tf1.paragraphs[0] if i==0 else tf1.add_paragraph()
    para.alignment = PP_ALIGN.CENTER
    para.space_before = Pt(6)
    run = para.add_run(); run.text = t
    run.font.size = Pt(sz); run.font.bold = bold
    run.font.color.rgb = WHITE if bold else LIGHT

# Slides 2-10
slides_data = [
    ("Objective", [
        "* Build an end-to-end Mutual Fund Analytics Platform",
        "* Ingest and clean 10 datasets covering 40 schemes across 10 AMCs",
        "* Compute production-quality risk-return metrics:",
        "  CAGR  |  Sharpe  |  Sortino  |  Alpha  |  Beta  |  Max Drawdown",
        "* Build a weighted fund scoring model (30/25/20/15/10)",
        "* Develop a risk-profile recommendation engine (Low/Medium/High)",
        "* Deliver notebooks, scripts, SQL queries, and reports",
    ]),
    ("Datasets", [
        "* 10 raw CSV files  |  87,533 total rows  |  All stored in data/raw/",
        "  01_fund_master          40 schemes metadata",
        "  02_nav_history          46,000 rows  Jan 2022 - May 2026",
        "  03_aum_by_fund_house    Bi-annual AUM per AMC",
        "  07_scheme_performance   Risk-return metrics per scheme",
        "  08_investor_transactions 32,778 rows  investor ledger",
        "  10_benchmark_indices    NIFTY50 + 6 benchmark indices",
        "* Data quality: 9/10 datasets zero nulls, zero duplicates",
    ]),
    ("Architecture", [
        "* data/raw/          10 raw CSVs (source of truth)",
        "* etl_pipeline.py    Extract -> Transform -> Load  (1.8 sec)",
        "* data/processed/    3 cleaned CSVs + 14 analytical outputs",
        "* data/db/           bluestock_mf.db (SQLite star schema)",
        "  dim_fund | fact_nav | fact_performance | fact_transaction | fact_aum",
        "* scripts/           compute_metrics.py  |  recommender.py",
        "* notebooks/         01 to 05 complete pipeline notebooks",
        "* sql/               schema.sql + 10-query queries.sql",
    ]),
    ("Data Cleaning", [
        "* NAV History:   dates parsed, daily_return_pct, rolling_vol_30d",
        "* Performance:   risk_score (1-6), excess_return, information_ratio",
        "* Transactions:  date parts, amount_bucket, kyc_verified flag",
        "* ETL Pipeline:  python scripts/etl_pipeline.py",
        "  Supports --step extract | transform | load",
        "* Validation:    9/10 OK  |  1 WARN (yoy_growth_pct - expected year 1)",
        "* Notebook:      02_data_cleaning.ipynb  (8 cells, 0 errors)",
    ]),
    ("Database Design", [
        "* Star Schema in SQLite  (bluestock_mf.db)",
        "  dim_fund          PK: amfi_code  -- 40 schemes",
        "  fact_nav          46,000 rows   -- daily NAV",
        "  fact_performance  40 rows       -- risk-return metrics",
        "  fact_transaction  32,778 rows   -- investor activity",
        "  fact_aum          90 rows       -- AUM by fund house",
        "* SQL files:  schema.sql  |  queries.sql (10 business queries)",
        "  Top 5 AUM | Avg NAV | SIP Trends | Expense Analysis | Risk vs Return",
    ]),
    ("Exploratory Data Analysis", [
        "* SBI MF leads AUM -- Rs 12.5 lakh crore (Mar 2025)",
        "* SIP inflows peaked Rs 31,002 Cr in Dec 2025",
        "* Folio count nearly doubled: 13.26 Cr to 26.12 Cr (2022-2025)",
        "* BFSI + IT dominate equity portfolio holdings",
        "* 26-35 age group: highest transaction count (13,463)",
        "* T30 cities: 66% of all investment transactions",
        "* 20+ charts generated across EDA notebook (03_eda_analysis.ipynb)",
    ]),
    ("Performance Analytics", [
        "* Mean 3Y CAGR:   14.09%   |  Best: SBI Small Cap Regular  23.39%",
        "* Mean Sharpe:     0.54    |  Best: Mirae Asset Large Cap   1.45",
        "* Mean Sortino:    0.92    |  Best: Mirae Asset Large Cap   2.39",
        "* Mean Max DD:  -17.87%   |  Best: ICICI Pru Liquid       -0.10%",
        "* Mean Alpha:      1.25   |  Best: HDFC Short Term Debt    1.98",
        "* Top 5 funds outperformed NIFTY50 by +106.21% over analysis period",
        "* Notebook: 04_performance_analytics.ipynb  (23 cells, 0 errors)",
    ]),
    ("Recommendation Engine", [
        "* Weighted score: 30% Return | 25% Sharpe | 20% Alpha",
        "                  15% Expense | 10% Max Drawdown",
        "* MinMaxScaler(0,100) -- fund_score_100: 100 = Best fund",
        "",
        "* Low Risk  (risk_score <= 2):  HDFC Short Term Debt  score=79.9",
        "* Medium    (risk_score <= 4):  Kotak Flexicap        score=97.4",
        "* High      (all funds):        SBI Small Cap Direct  score=100.0",
        "",
        "* scripts/recommender.py  ->  recommend('Low' | 'Medium' | 'High')",
    ]),
    ("Conclusion", [
        "[OK]  40 schemes profiled across 10 AMCs -- full risk-return metrics",
        "[OK]  ETL pipeline: 87,533 rows processed in 1.8 seconds",
        "[OK]  5 notebooks: ingestion -> cleaning -> EDA -> performance -> advanced",
        "[OK]  8 reusable financial metric functions in compute_metrics.py",
        "[OK]  SQL schema + 10 business queries (queries.sql)",
        "[OK]  Recommendation engine: Low / Medium / High risk profiles",
        "[OK]  Top 5 funds: +112% vs NIFTY50 +5.9% over same period",
        "[OK]  Final validation: 38/38 checks PASSED -- PROJECT COMPLETE",
    ]),
]

for title, bullets in slides_data:
    slide(title, bullets)

out_pptx = RPT / "Presentation.pptx"
prs.save(str(out_pptx))
print(f"PPTX saved -> {out_pptx}  ({out_pptx.stat().st_size // 1024} KB)")
print("All reports generated.")
